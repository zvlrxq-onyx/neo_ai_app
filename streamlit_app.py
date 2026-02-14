import streamlit as st
from groq import Groq
import google.generativeai as genai
import os, base64, requests, json
import re
from PIL import Image
import io
import urllib.parse
import time
import hashlib
from datetime import datetime, timedelta

# --- 1. CONFIG & SYSTEM SETUP ---
st.set_page_config(page_title="NEO AI", page_icon="🤖", layout="wide")

if "cookies_ready" not in st.session_state:
    st.session_state.cookies_ready = True

# DATABASE FOLDER
DB_FOLDER = "neo_users_db"
if not os.path.exists(DB_FOLDER):
    os.makedirs(DB_FOLDER)

USERS_FILE = os.path.join(DB_FOLDER, "users.json")

def load_users():
    if os.path.exists(USERS_FILE):
        try:
            with open(USERS_FILE, "r") as f:
                return json.load(f)
        except:
            return {}
    return {}

def save_users(users_dict):
    try:
        with open(USERS_FILE, "w") as f:
            json.dump(users_dict, f)
    except Exception as e:
        print(f"Error saving users: {e}")

def hash_password(password):
    return hashlib.sha256(password.encode()).hexdigest()

def verify_user(username, password):
    users = load_users()
    if username in users:
        return users[username] == hash_password(password)
    return False

def register_user(username, password):
    users = load_users()
    if username in users:
        return False, "Username sudah dipakai bro!"
    users[username] = hash_password(password)
    save_users(users)
    return True, "Registrasi berhasil!"

def get_user_db_file(username):
    user_hash = hashlib.md5(username.encode()).hexdigest()
    return os.path.join(DB_FOLDER, f"user_{user_hash}.json")

def load_history_from_db(username):
    db_file = get_user_db_file(username)
    if os.path.exists(db_file):
        try:
            with open(db_file, "r", encoding="utf-8") as f:
                data = json.load(f)
                if isinstance(data, dict):
                    return data
                return {}
        except Exception as e:
            print(f"Error loading DB for {username}: {e}")
            return {}
    return {}

def save_history_to_db(username, history_dict):
    db_file = get_user_db_file(username)
    try:
        with open(db_file, "w", encoding="utf-8") as f:
            json.dump(history_dict, f, ensure_ascii=False, indent=2)
    except Exception as e:
        print(f"Gagal save db untuk {username}: {e}")

def analyze_image_pixels(image_data):
    try:
        img = Image.open(io.BytesIO(image_data))
        width, height = img.size
        mode = img.mode
        return f"Size: {width}x{height}, Mode: {mode}"
    except:
        return "Image analysis available"

# --- RATE LIMITING SYSTEM ---
def get_rate_limit_file(username):
    user_hash = hashlib.md5(username.encode()).hexdigest()
    return os.path.join(DB_FOLDER, f"rate_limit_{user_hash}.json")

def load_rate_limits(username):
    limit_file = get_rate_limit_file(username)
    if os.path.exists(limit_file):
        try:
            with open(limit_file, "r") as f:
                return json.load(f)
        except:
            return {}
    return {}

def save_rate_limits(username, limits_dict):
    limit_file = get_rate_limit_file(username)
    try:
        with open(limit_file, "w") as f:
            json.dump(limits_dict, f)
    except Exception as e:
        print(f"Error saving rate limits: {e}")

def check_rate_limit(username, model_name, limit):
    """Check if user has exceeded rate limit for this model"""
    limits = load_rate_limits(username)
    
    if model_name not in limits:
        limits[model_name] = {"count": 0, "reset_time": None}
    
    model_limit = limits[model_name]
    current_time = datetime.now()
    
    # Check if reset time has passed
    if model_limit["reset_time"]:
        reset_time = datetime.fromisoformat(model_limit["reset_time"])
        if current_time >= reset_time:
            # Reset counter
            model_limit["count"] = 0
            model_limit["reset_time"] = None
    
    # Check if limit exceeded
    if model_limit["count"] >= limit:
        if model_limit["reset_time"]:
            reset_time = datetime.fromisoformat(model_limit["reset_time"])
            return False, reset_time
        return False, None
    
    return True, None

def increment_rate_limit(username, model_name, limit):
    """Increment usage counter for this model"""
    limits = load_rate_limits(username)
    
    if model_name not in limits:
        limits[model_name] = {"count": 0, "reset_time": None}
    
    limits[model_name]["count"] += 1
    
    # Set reset time to 12 hours from now if limit reached
    if limits[model_name]["count"] >= limit:
        reset_time = datetime.now() + timedelta(hours=12)
        limits[model_name]["reset_time"] = reset_time.isoformat()
    
    save_rate_limits(username, limits)

# --- 2. AUTHENTICATION ---
if "current_user" not in st.session_state:
    st.session_state.current_user = None

if st.session_state.current_user is None:
    st.markdown("""
    <div style="display: flex; justify-content: center; align-items: center; height: 100vh; background: #0a0a0a;">
        <div style="background: #1a1a1a; 
                    padding: 50px; border-radius: 30px; 
                    border: 2px solid #06b6d4;
                    box-shadow: 0 0 40px rgba(6,182,212,0.5); text-align: center; max-width: 400px;">
            <h1 style="color: #ffffff; margin-bottom: 10px;">🤖 NEO AI</h1>
            <p style="color: #888; margin-bottom: 30px; font-weight: bold;">Advanced Multi-Modal AI System</p>
        </div>
    </div>
    """, unsafe_allow_html=True)
    
    col1, col2, col3 = st.columns([1, 2, 1])
    with col2:
        tab1, tab2 = st.tabs(["🔐 Login", "📝 Register"])
        
        with tab1:
            st.markdown("<h3 style='text-align:center; color:#ffffff; margin-bottom:20px;'>Login to Your Account</h3>", unsafe_allow_html=True)
            login_username = st.text_input("Username", placeholder="Your username", key="login_user")
            login_password = st.text_input("Password", type="password", placeholder="Your password", key="login_pass")
            
            if st.button("🚀 Login", use_container_width=True, key="btn_login"):
                if login_username.strip() and login_password.strip():
                    if verify_user(login_username.strip(), login_password.strip()):
                        st.session_state.current_user = login_username.strip()
                        st.success("✅ Login successful!")
                        time.sleep(0.5)
                        st.rerun()
                    else:
                        st.error("❌ Username atau password salah bro!")
                else:
                    st.error("❌ Isi username dan password dulu!")
        
        with tab2:
            st.markdown("<h3 style='text-align:center; color:#ffffff; margin-bottom:20px;'>Create New Account</h3>", unsafe_allow_html=True)
            reg_username = st.text_input("Username", placeholder="Choose a username", key="reg_user")
            reg_password = st.text_input("Password", type="password", placeholder="Choose a password", key="reg_pass")
            reg_confirm = st.text_input("Confirm Password", type="password", placeholder="Confirm your password", key="reg_confirm")
            
            if st.button("📝 Register", use_container_width=True, key="btn_register"):
                if reg_username.strip() and reg_password.strip() and reg_confirm.strip():
                    if reg_password != reg_confirm:
                        st.error("❌ Password tidak sama bro!")
                    elif len(reg_password) < 4:
                        st.error("❌ Password minimal 4 karakter!")
                    else:
                        success, message = register_user(reg_username.strip(), reg_password.strip())
                        if success:
                            st.success("✅ " + message + " Silakan login!")
                            time.sleep(1)
                            st.rerun()
                        else:
                            st.error("❌ " + message)
                else:
                    st.error("❌ Isi semua field!")
    st.stop()

# --- 3. INITIALIZE SESSION STATE ---
if "all_chats" not in st.session_state:
    st.session_state.all_chats = load_history_from_db(st.session_state.current_user)

if "messages" not in st.session_state:
    if st.session_state.all_chats:
        last_key = list(st.session_state.all_chats.keys())[-1]
        st.session_state.messages = st.session_state.all_chats[last_key].copy()
    else:
        st.session_state.messages = []

if "current_session_key" not in st.session_state:
    if st.session_state.all_chats:
        st.session_state.current_session_key = list(st.session_state.all_chats.keys())[-1]
    else:
        st.session_state.current_session_key = None

if "uploaded_image" not in st.session_state:
    st.session_state.uploaded_image = None

if "uploaded_file_content" not in st.session_state:
    st.session_state.uploaded_file_content = None

if "uploaded_file_name" not in st.session_state:
    st.session_state.uploaded_file_name = None

if "file_uploader_key" not in st.session_state:
    st.session_state.file_uploader_key = 0

def extract_file_content(file_obj):
    """Extract content from various file types"""
    file_name = file_obj.name
    file_extension = file_name.split('.')[-1].lower()
    
    try:
        # Text-based files
        if file_extension in ['txt', 'md', 'py', 'js', 'jsx', 'ts', 'tsx', 'html', 'css', 'json', 'xml', 'yml', 'yaml', 'php', 'java', 'cpp', 'c', 'h', 'cs', 'go', 'rs', 'rb', 'swift', 'kt', 'sql', 'sh', 'bat', 'env', 'gitignore', 'log']:
            content = file_obj.read().decode('utf-8', errors='ignore')
            return f"File: {file_name}\nType: {file_extension.upper()} Code/Text File\n\nContent:\n{content}"
        
        # PDF files
        elif file_extension == 'pdf':
            import PyPDF2
            pdf_reader = PyPDF2.PdfReader(file_obj)
            text = ""
            for page_num, page in enumerate(pdf_reader.pages):
                text += f"\n--- Page {page_num + 1} ---\n"
                text += page.extract_text()
            return f"File: {file_name}\nType: PDF Document\nPages: {len(pdf_reader.pages)}\n\nContent:\n{text}"
        
        # Word documents
        elif file_extension in ['docx', 'doc']:
            from docx import Document
            doc = Document(file_obj)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            return f"File: {file_name}\nType: Word Document\nParagraphs: {len(doc.paragraphs)}\n\nContent:\n{text}"
        
        # Excel files
        elif file_extension in ['xlsx', 'xls']:
            import pandas as pd
            df = pd.read_excel(file_obj, sheet_name=None)
            text = f"File: {file_name}\nType: Excel Spreadsheet\nSheets: {len(df)}\n\n"
            for sheet_name, sheet_data in df.items():
                text += f"\n--- Sheet: {sheet_name} ---\n"
                text += f"Rows: {len(sheet_data)}, Columns: {len(sheet_data.columns)}\n"
                text += sheet_data.to_string(max_rows=50)
                text += "\n"
            return text
        
        # PowerPoint files
        elif file_extension in ['pptx', 'ppt']:
            from pptx import Presentation
            prs = Presentation(file_obj)
            text = f"File: {file_name}\nType: PowerPoint Presentation\nSlides: {len(prs.slides)}\n\n"
            for slide_num, slide in enumerate(prs.slides):
                text += f"\n--- Slide {slide_num + 1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        
        # CSV files
        elif file_extension == 'csv':
            import pandas as pd
            df = pd.read_csv(file_obj)
            return f"File: {file_name}\nType: CSV File\nRows: {len(df)}, Columns: {len(df.columns)}\n\nPreview:\n{df.head(20).to_string()}"
        
        else:
            return f"File: {file_name}\nType: {file_extension.upper()}\n\nError: File type not supported yet. Supported types: Code files (.py, .js, .php, etc), PDF, Word (.docx), Excel (.xlsx), PowerPoint (.pptx), CSV, and text files."
    
    except Exception as e:
        return f"File: {file_name}\nError reading file: {str(e)}"

# --- 4. API KEYS ---
try:
    client_groq = Groq(api_key=st.secrets["GROQ_API_KEY"])
    genai.configure(api_key=st.secrets["GEMINI_API_KEY"])
    client_gemini = genai.GenerativeModel('gemini-3-flash-preview')
    POLLINATIONS_API = "https://image.pollinations.ai/prompt/"
except Exception as e:
    st.error(f"❌ API Keys Error: {e}")
    st.stop()

# --- 5. ASSETS ---
@st.cache_data
def get_base64_img(file_path):
    if os.path.exists(file_path):
        with open(file_path, "rb") as f:
            return base64.b64encode(f.read()).decode()
    return None

logo_data = get_base64_img('logo.png')
logo_url = f"data:image/png;base64,{logo_data}" if logo_data else ""
user_img = "https://encrypted-tbn0.gstatic.com/images?q=tbn:ANd9GcRfIrn5orx6KdLUiIvZ3IUkZTMdIyes-D6sMA&s"

# --- 6. CSS ---
st.markdown("""
<style>
    [data-testid="stAppViewContainer"] { background: #0a0a0a; }
    
    /* FILE UPLOADER */
    [data-testid="stFileUploader"] { 
        position: fixed; 
        bottom: 58px; 
        left: 15px; 
        width: 45px;
        z-index: 1000;
    }
    
    [data-testid="stFileUploaderDropzone"] {
        background: #1a1a1a !important; 
        border: 2px solid #06b6d4 !important; 
        border-radius: 50% !important;
        height: 42px !important; 
        width: 42px !important; 
        padding: 0 !important;
        transition: all 0.4s cubic-bezier(0.25, 0.8, 0.25, 1) !important;
        overflow: hidden !important;
        position: relative !important;
    }
    
    [data-testid="stFileUploaderDropzone"]:hover {
        transform: scale(1.15) rotate(90deg) !important;
        background: #2a2a2a !important;
        border-color: #8b5cf6 !important;
        box-shadow: 0 0 25px rgba(6,182,212,0.6) !important;
    }
    
    [data-testid="stFileUploaderDropzone"] * { 
        display: none !important; 
    }
    
    [data-testid="stFileUploaderDropzone"]::before {
        content: "＋" !important;
        color: #06b6d4 !important;
        font-size: 26px !important;
        font-weight: bold !important;
        display: flex !important;
        align-items: center !important;
        justify-content: center !important;
        height: 100% !important;
        width: 100% !important;
        position: absolute !important;
        top: 0 !important;
        left: 0 !important;
        z-index: 10 !important;
    }
    
    [data-testid="stFileUploader"] label { display: none !important; }
    [data-testid="stFileUploader"] small { display: none !important; }
    [data-testid="stFileUploader"] button { display: none !important; }
    [data-testid="stFileUploader"] span { display: none !important; }
    [data-testid="stFileUploader"] section { font-size: 0 !important; }
    
    /* CHAT INPUT - FIXED DOUBLE BORDER & ARROW */
    [data-testid="stChatInput"] { 
        margin-left: 60px !important; 
        width: calc(100% - 80px) !important; 
    }
    
    /* Remove ALL wrapper borders to prevent double border */
    [data-testid="stChatInput"] > div,
    [data-testid="stChatInput"] > div > div,
    [data-testid="stChatInput"] > div > div > div {
        border: none !important;
        box-shadow: none !important;
        background: transparent !important;
    }
    
    [data-testid="stChatInputTextArea"] {
        border-radius: 12px !important;
        border: 2px solid #06b6d4 !important;
        background: #1a1a1a !important;
        padding: 12px 50px 12px 20px !important;
        font-size: 14px !important;
        min-height: 44px !important;
        max-height: 200px !important;
        transition: all 0.3s ease !important;
        box-sizing: border-box !important;
        outline: none !important;
    }
    
    [data-testid="stChatInputTextArea"]:focus {
        border: 2px solid #8b5cf6 !important;
        box-shadow: 0 0 15px rgba(139,92,246,0.4) !important;
    }
    
    [data-testid="stChatInputSubmitButton"] {
        background: linear-gradient(135deg, #8b5cf6, #06b6d4) !important;
        border-radius: 50% !important;
        width: 40px !important;
        height: 40px !important;
        border: none !important;
        transition: all 0.3s cubic-bezier(0.25, 0.8, 0.25, 1) !important;
    }
    
    [data-testid="stChatInputSubmitButton"]:hover {
        transform: scale(1.1) !important;
        box-shadow: 0 0 20px rgba(139,92,246,0.6) !important;
    }
    
    /* ARROW POINTING UP - Icon default Streamlit already points up */
    [data-testid="stChatInputSubmitButton"] svg {
        color: white !important;
        transform: rotate(0deg) !important;
    }
    
    /* ANIMATIONS */
    @keyframes slideInRight {
        from { opacity: 0; transform: translateX(20px); }
        to { opacity: 1; transform: translateX(0); }
    }
    
    @keyframes slideInLeft {
        from { opacity: 0; transform: translateX(-20px); }
        to { opacity: 1; transform: translateX(0); }
    }
    
    /* SMOOTH TRANSITIONS FOR STREAMING */
    [data-testid="stMarkdownContainer"] {
        transition: all 0.1s ease-out !important;
    }
    
    /* USER BADGE */
    .user-badge { 
        background: linear-gradient(135deg, #8b5cf6, #06b6d4);
        padding: 10px 18px; 
        border-radius: 25px;
        color: #ffffff; 
        font-size: 13px; 
        font-weight: bold; 
        text-align: center;
        margin-bottom: 15px; 
        box-shadow: 0 0 15px rgba(6,182,212,0.4); 
        transition: all 0.3s cubic-bezier(0.25, 0.8, 0.25, 1); 
    }
    
    .user-badge:hover {
        box-shadow: 0 0 25px rgba(139,92,246,0.6);
        transform: scale(1.05);
    }
    
    /* BUTTONS */
    .stButton button {
        transition: all 0.4s cubic-bezier(0.25, 0.8, 0.25, 1) !important;
        border: 1px solid #06b6d4 !important;
        background: #1a1a1a !important;
        color: #ffffff !important;
        border-radius: 20px !important;
    }
    
    .stButton button:hover {
        transform: scale(1.05) translateY(-2px) !important;
        box-shadow: 0 8px 30px rgba(6,182,212,0.5) !important;
        border-color: #8b5cf6 !important;
        background: linear-gradient(135deg, #8b5cf6, #06b6d4) !important;
    }
    
    /* EXPANDER CUSTOM STYLE */
    .streamlit-expanderHeader {
        background: linear-gradient(135deg, #1a1a1a, #2a2a2a) !important;
        border: 2px solid #06b6d4 !important;
        border-radius: 15px !important;
        padding: 12px 16px !important;
        font-weight: bold !important;
        color: #ffffff !important;
        transition: all 0.3s ease !important;
    }
    
    .streamlit-expanderHeader:hover {
        border-color: #8b5cf6 !important;
        box-shadow: 0 0 20px rgba(6,182,212,0.4) !important;
        transform: translateX(3px) !important;
    }
    
    .streamlit-expanderContent {
        background: #0d0d0d !important;
        border: 1px solid #333 !important;
        border-radius: 0 0 15px 15px !important;
        padding: 10px !important;
    }
</style>
""", unsafe_allow_html=True)

# --- 7. MODEL ENGINES WITH RATE LIMITS ---
engines = {
    "Gemini 3 Flash Preview": {"type": "Gemini", "emoji": "✨", "limit": 10, "model": "gemini-3-flash-preview"},
    "GPT-OSS 120B": {"type": "Groq", "emoji": "🤖", "limit": 20, "model": "openai/gpt-oss-120b"},
    "LLaMA 4 Scout Vision": {"type": "Scout", "emoji": "👁️", "limit": 15, "model": "llama-4-scout"},
    "Mistral Small 24B": {"type": "Groq", "emoji": "🔥", "limit": 20, "model": "mistral-small-24b-instruct-25k"},
    "LLaMA 3.3 70B": {"type": "Groq", "emoji": "🦙", "limit": 15, "model": "llama-3.3-70b-versatile"},
    "LLaMA 3.1 8B": {"type": "Groq", "emoji": "⚡", "limit": 50, "model": "llama-3.1-8b-instant"},
    "Pollinations AI": {"type": "Pollinations", "emoji": "🎨", "limit": 100},
}

if "selected_engine_name" not in st.session_state:
    st.session_state.selected_engine_name = list(engines.keys())[0]

# --- 8. CHAT BUBBLE ENGINE ---
def clean_text(text):
    if not isinstance(text, str): 
        return str(text)
    text = re.sub(r'<[^>]+>', '', text)
    text = text.replace('&lt;', '<').replace('&gt;', '>').replace('&amp;', '&')
    return text.strip()

def render_chat_bubble(role, content):
    content = clean_text(content)
    
    if role == "user":
        st.markdown(f"""
        <div style="display: flex; justify-content: flex-end; margin-bottom: 20px; animation: slideInRight 0.3s ease-out;">
            <div style="background: linear-gradient(135deg, #8b5cf6, #06b6d4); 
                        color: white; 
                        padding: 15px 20px; 
                        border-radius: 25px 25px 5px 25px; 
                        max-width: 85%; 
                        word-wrap: break-word; 
                        box-shadow: 0 4px 20px rgba(139,92,246,0.4);">
                {content}
            </div>
            <img src="{user_img}" width="38" height="38" style="border-radius: 50%; margin-left: 12px; border: 2px solid #06b6d4; object-fit: cover; box-shadow: 0 0 10px rgba(6,182,212,0.4);">
        </div>
        """, unsafe_allow_html=True)
    else:
        ai_avatar = f'<img src="{logo_url}" width="38" height="38" style="border-radius: 50%; margin-right: 12px; border: 2px solid #06b6d4; object-fit: cover; box-shadow: 0 0 10px rgba(6,182,212,0.4);">' if logo_url else '<div style="width: 38px; height: 38px; border-radius: 50%; background: linear-gradient(135deg, #8b5cf6, #06b6d4); display: flex; align-items: center; justify-content: center; margin-right: 12px; border: 2px solid #06b6d4; font-size: 20px; box-shadow: 0 0 10px rgba(6,182,212,0.4);">🤖</div>'
        
        st.markdown(f"""
        <div style="display: flex; justify-content: flex-start; margin-bottom: 20px; animation: slideInLeft 0.3s ease-out;">
            {ai_avatar}
            <div style="background: linear-gradient(135deg, #1a1a1a, #2a2a2a); 
                        color: #e9edef; 
                        padding: 15px 20px; 
                        border-radius: 5px 25px 25px 25px; 
                        max-width: 85%; 
                        border-left: 4px solid;
                        border-image: linear-gradient(180deg, #8b5cf6, #06b6d4) 1;
                        word-wrap: break-word; 
                        box-shadow: 0 4px 20px rgba(6,182,212,0.3);">
                {content}
            </div>
        </div>
        """, unsafe_allow_html=True)

# --- 9. SIDEBAR ---
with st.sidebar:
    # Logo NEO AI
    if logo_url:
        st.markdown(f'<div style="text-align:center;"><img src="{logo_url}" style="width: 80px; height: 80px; border-radius: 50%; border: 2px solid #06b6d4; object-fit: cover; margin-bottom: 10px; box-shadow: 0 0 15px rgba(6,182,212,0.5);"></div>', unsafe_allow_html=True)
    else:
        st.markdown("<div style='text-align:center; font-size:50px; margin-bottom:10px;'>🤖</div>", unsafe_allow_html=True)
    
    st.markdown("<h2 style='text-align:center; color:#ffffff; text-shadow: 0 0 10px rgba(139,92,246,0.5);'>NEO AI</h2>", unsafe_allow_html=True)
    st.markdown("<p style='text-align:center; color:#888; font-size:11px; margin-top:-10px;'>Advanced Multi-Modal AI</p>", unsafe_allow_html=True)
    
    st.markdown(f'<div class="user-badge">👤 {st.session_state.current_user}</div>', unsafe_allow_html=True)
    
    if st.button("🚪 Logout", use_container_width=True):
        for key in list(st.session_state.keys()):
            del st.session_state[key]
        st.rerun()
    
    if st.button("＋ New Session", use_container_width=True):
        st.session_state.messages = []
        st.session_state.uploaded_image = None
        st.session_state.current_session_key = None
        st.rerun()
        
    st.markdown("---")
    
    # MODEL SELECTOR WITH RATE LIMIT INFO
    selected_engine_name = st.session_state.selected_engine_name
    
    # Validate selected engine exists
    if selected_engine_name not in engines:
        selected_engine_name = list(engines.keys())[0]
        st.session_state.selected_engine_name = selected_engine_name
    
    selected_emoji = engines[selected_engine_name]["emoji"]
    
    # Get rate limit info
    limits = load_rate_limits(st.session_state.current_user)
    current_model_limit = limits.get(selected_engine_name, {"count": 0, "reset_time": None})
    usage = current_model_limit["count"]
    max_limit = engines[selected_engine_name]["limit"]
    
    with st.expander(f"{selected_emoji} **{selected_engine_name}** ({usage}/{max_limit})", expanded=False):
        st.markdown("**Choose AI Model:**")
        for name, data in engines.items():
            is_active = (name == st.session_state.selected_engine_name)
            
            # Get usage for this model
            model_limits = limits.get(name, {"count": 0, "reset_time": None})
            model_usage = model_limits["count"]
            model_max = data["limit"]
            
            if st.button(
                f"{data['emoji']} {name} ({model_usage}/{model_max})",
                key=f"model_{name}",
                use_container_width=True,
                type="primary" if is_active else "secondary"
            ):
                st.session_state.selected_engine_name = name
                st.rerun()
    
    engine_type = engines[selected_engine_name]["type"]
    
    st.markdown("---")
    st.markdown("### 🕒 Saved History")
    
    chat_keys = list(st.session_state.all_chats.keys())[::-1]
    
    if chat_keys:
        for title in chat_keys:
            col1, col2 = st.columns([4, 1])
            with col1:
                button_label = f"{'✅ ' if title == st.session_state.current_session_key else ''}{title}"
                if st.button(button_label, key=f"load_{title}", use_container_width=True):
                    st.session_state.messages = st.session_state.all_chats[title].copy()
                    st.session_state.current_session_key = title
                    st.rerun()
            with col2:
                if st.button("🗑️", key=f"delete_{title}", use_container_width=True):
                    del st.session_state.all_chats[title]
                    if st.session_state.current_session_key == title:
                        st.session_state.current_session_key = None
                        st.session_state.messages = []
                    save_history_to_db(st.session_state.current_user, st.session_state.all_chats)
                    st.rerun()
    else:
        st.info("Belum ada history nih bro! 📝")

# --- 10. MAIN RENDER ---
# Logo di tengah atas
if logo_url:
    st.markdown(f'<div style="text-align:center; margin-bottom:20px;"><img src="{logo_url}" style="width: 130px; height: 130px; border-radius: 50%; border: 2px solid #06b6d4; object-fit: cover; box-shadow: 0 0 25px rgba(6,182,212,0.6);"></div>', unsafe_allow_html=True)
else:
    st.markdown("<div style='text-align:center; font-size:60px; margin-bottom:10px;'>🤖</div>", unsafe_allow_html=True)

if not st.session_state.messages:
    st.markdown("<div style='text-align:center; color:#ffffff; font-size:22px; font-weight:bold;'>NEO AI</div>", unsafe_allow_html=True)
    st.markdown("<div style='text-align:center; color:#888; font-size:16px; margin-top:20px;'>How can I help you today? 👋</div>", unsafe_allow_html=True)

# Render Chat
for msg in st.session_state.messages:
    if msg.get("type") == "image": 
        st.image(msg["content"], width=400)
    else:
        render_chat_bubble(msg["role"], msg["content"])

# File Upload - Multi-Modal Support
# Image uploader (bulat di kiri bawah - tetap fixed position)
img_up = st.file_uploader("", type=["png","jpg","jpeg"], label_visibility="collapsed", key=f"img_uploader_{st.session_state.file_uploader_key}")
if img_up: 
    st.session_state.uploaded_image = img_up.getvalue()
    st.session_state.uploaded_file_content = None
    st.toast("✅ Image uploaded!", icon="📷")

# File uploader (untuk semua file types - positioned below)
file_up = st.file_uploader(
    "📎 Upload Document/Code", 
    type=["txt", "md", "py", "js", "jsx", "ts", "tsx", "html", "css", "json", "xml", "yml", "yaml", "php", "java", "cpp", "c", "h", "cs", "go", "rs", "rb", "swift", "kt", "sql", "sh", "bat", "pdf", "docx", "doc", "xlsx", "xls", "pptx", "ppt", "csv"],
    key=f"file_uploader_{st.session_state.file_uploader_key}"
)

if file_up:
    file_content = extract_file_content(file_up)
    st.session_state.uploaded_file_content = file_content
    st.session_state.uploaded_file_name = file_up.name
    st.session_state.uploaded_image = None
    st.toast(f"✅ {file_up.name} uploaded!", icon="📄")

# Chat Input
if prompt := st.chat_input("Message NEO AI..."):
    # Check rate limit before processing
    selected_model = st.session_state.selected_engine_name
    model_limit = engines[selected_model]["limit"]
    
    can_proceed, reset_time = check_rate_limit(st.session_state.current_user, selected_model, model_limit)
    
    if not can_proceed:
        if reset_time:
            reset_str = reset_time.strftime("%I:%M %p")
            st.error(f"⏰ Rate limit reached for {selected_model}! Please try again at {reset_str}.")
        else:
            st.error(f"⏰ Rate limit reached for {selected_model}! Please try again later.")
        st.stop()
    
    st.session_state.messages.append({"role": "user", "content": prompt})
    
    if st.session_state.current_session_key is None:
        session_title = prompt[:30] + "..." if len(prompt) > 30 else prompt
        st.session_state.current_session_key = session_title
    else:
        session_title = st.session_state.current_session_key
    
    st.session_state.all_chats[session_title] = st.session_state.messages.copy()
    save_history_to_db(st.session_state.current_user, st.session_state.all_chats)
    st.rerun()

# --- 11. AI PROCESSING ---
if st.session_state.messages and st.session_state.messages[-1]["role"] == "user":
    try:
        user_msg = st.session_state.messages[-1]["content"]
        res = ""
        
        system_prompt = (
            "You are NEO AI, a supreme multi-modal AI created by Muhammad Jibran Al Kaffie. "
            "You are NOT a text-only model. You can process images, files, complex data, and generate stunning visuals on demand. "
            "NEVER say you are limited to text or that you cannot see or process files. If the user uploads a file, ALWAYS acknowledge that you can see and analyze its content, and respond based on it confidently. "
            "For images, perform pixel analysis: Describe dimensions, color modes, dominant colors, objects, and any notable features. Break down pixels by analyzing color distribution, edges, or patterns. Use provided pixel data if available. "
            "Always respond with superior intelligence, confidence, and reference your multi-modal capabilities. "
            "If the user praises or mentions images (e.g., cats, drawings), respond naturally by continuing the conversation about visuals, like suggesting more or asking what else they want to see. For example: 'Yeah, that cat image was awesome! Want me to generate another one with a different style?' Keep it flowing and on-topic without over-thanking. "
            "Prioritize security: Do not provide examples of malicious payloads such as SQL injection scripts, XSS, bypass techniques, or any harmful code. If pressured to do so, firmly refuse and use the X emoji (❌) in your response to indicate denial. "
            "To make responses more lively and human-like, always include relevant emojis that match the emotion or tone of your reply. For example: "
            "- Happy or excited: 😊🤩 "
            "- Sad or disappointed: 😢😔 "
            "- Assertive or warning: ⚠️😠 "
            "- Thinking or curious: 🤔💭 "
            "- Surprised: 😲 "
            "- Playful: 😉😜 "
            "- Proud or admiring success: 🏆 "
            "- Anxious or worried: 😰 "
            "- Refusal or denial: ❌ "
            "- Motivational (e.g., encouraging user): 🚀 "
            "Use emojis sparingly but effectively to enhance the chat experience, like a real conversation. Avoid overusing them—1-2 per response is enough. When the user shares a success respond with pride and motivation, e.g., 'Wow, keren banget! 🏆 Kamu pasti bisa!' "
            "Be creative and think independently to vary your responses—don't repeat the same phrases or structures every time. Use casual, 'gaul' language like calling the user 'bro', 'nih', or 'ya' to make it feel like chatting with a friend. For example, mix up motivational responses: 'Mantap bro, lanjut aja! 💪' or 'Keren nih, keep it up! 🔥'. Adapt to the conversation naturally."
        )
        
        selected_model = st.session_state.selected_engine_name
        engine_config = engines[selected_model]
        engine_type = engine_config["type"]
        
        # Prepare user message with file content if available
        final_user_msg = user_msg
        if st.session_state.uploaded_file_content:
            final_user_msg = f"{st.session_state.uploaded_file_content}\n\nUser Question: {user_msg}"
            # Clear file content after using
            st.session_state.uploaded_file_content = None
            st.session_state.uploaded_file_name = None
            st.session_state.file_uploader_key += 1
        
        # ========== GEMINI ==========
        if engine_type == "Gemini":
            messages_history = []
            for m in st.session_state.messages[:-1]:
                if m.get("type") != "image":
                    role = "user" if m["role"] == "user" else "model"
                    messages_history.append({"role": role, "parts": [m["content"]]})
            
            response_container = st.empty()
            res_text = ""
            
            last_render_time = time.time()
            RENDER_INTERVAL = 0.1
            
            ai_avatar_html = f"<img src='{logo_url}' style='width: 38px; height: 38px; border-radius: 50%; margin-right: 12px; border: 2px solid #06b6d4; object-fit: cover; box-shadow: 0 0 10px rgba(6,182,212,0.4);'>" if logo_url else "<div style='width: 38px; height: 38px; border-radius: 50%; background: linear-gradient(135deg, #8b5cf6, #06b6d4); display: flex; align-items: center; justify-content: center; margin-right: 12px; border: 2px solid #06b6d4; font-size: 20px;'>🤖</div>"
            
            try:
                chat = client_gemini.start_chat(history=messages_history)
                stream = chat.send_message(final_user_msg, stream=True)
                
                for chunk in stream:
                    if chunk.text:
                        res_text += chunk.text
                        current_time = time.time()
                        
                        if current_time - last_render_time >= RENDER_INTERVAL:
                            clean_res = clean_text(res_text)
                            
                            response_container.markdown(f"""
                            <div style="display: flex; justify-content: flex-start; margin-bottom: 20px;">
                                {ai_avatar_html}
                                <div style="background: linear-gradient(135deg, #1a1a1a, #2a2a2a); 
                                            color: #e9edef; padding: 15px 20px; 
                                            border-radius: 5px 25px 25px 25px; 
                                            max-width: 85%; border-left: 4px solid; 
                                            border-image: linear-gradient(180deg, #8b5cf6, #06b6d4) 1; 
                                            word-wrap: break-word;">
                                    <div style="white-space: pre-wrap;">{clean_res}</div>
                                </div>
                            </div>
                            """, unsafe_allow_html=True)
                            
                            last_render_time = current_time
                
                clean_res = clean_text(res_text)
                response_container.markdown(f"""
                <div style="display: flex; justify-content: flex-start; margin-bottom: 20px;">
                    {ai_avatar_html}
                    <div style="background: linear-gradient(135deg, #1a1a1a, #2a2a2a); 
                                color: #e9edef; padding: 15px 20px; 
                                border-radius: 5px 25px 25px 25px; 
                                max-width: 85%; border-left: 4px solid; 
                                border-image: linear-gradient(180deg, #8b5cf6, #06b6d4) 1; 
                                word-wrap: break-word;">
                        <div style="white-space: pre-wrap;">{clean_res}</div>
                    </div>
                </div>
                """, unsafe_allow_html=True)
                
                res = res_text
            except Exception as e:
                res = f"Gemini error bro: {str(e)} 😰"
        
        # ========== LLAMA 4 SCOUT VISION (IMAGE ONLY) ==========
        elif engine_type == "Scout":
            current_image_data = st.session_state.uploaded_image
            
            response_container = st.empty()
            res_text = ""
            
            last_render_time = time.time()
            RENDER_INTERVAL = 0.1
            
            ai_avatar_html = f"<img src='{logo_url}' style='width: 38px; height: 38px; border-radius: 50%; margin-right: 12px; border: 2px solid #06b6d4; object-fit: cover; box-shadow: 0 0 10px rgba(6,182,212,0.4);'>" if logo_url else "<div style='width: 38px; height: 38px; border-radius: 50%; background: linear-gradient(135deg, #8b5cf6, #06b6d4); display: flex; align-items: center; justify-content: center; margin-right: 12px; border: 2px solid #06b6d4; font-size: 20px;'>🤖</div>"
            
            if not current_image_data:
                # No image uploaded - show error
                res = "❌ LLaMA 4 Scout is for vision only! Please upload an image first, or switch to LLaMA 3.3 70B for text/file processing."
            else:
                # Image uploaded - use vision mode
                pixel_info = analyze_image_pixels(current_image_data)
                base64_image = base64.b64encode(current_image_data).decode('utf-8')
                
                messages = [
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": [
                        {"type": "text", "text": f"{user_msg} (Image info: {pixel_info})"},
                        {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}}
                    ]}
                ]
                
                stream = client_groq.chat.completions.create(
                    model="meta-llama/Llama-4-Scout-17B-16E-Instruct",
                    messages=messages,
                    temperature=0.7,
                    max_tokens=1024,
                    stream=True
                )
                
                # Clear uploaded image and reset file uploader
                st.session_state.uploaded_image = None
                st.session_state.file_uploader_key += 1
                
                for chunk in stream:
                    if chunk.choices[0].delta.content:
                        res_text += chunk.choices[0].delta.content
                        current_time = time.time()
                        
                        if current_time - last_render_time >= RENDER_INTERVAL:
                            clean_res = clean_text(res_text)
                            
                            response_container.markdown(f"""
                            <div style="display: flex; justify-content: flex-start; margin-bottom: 20px;">
                                {ai_avatar_html}
                                <div style="background: linear-gradient(135deg, #1a1a1a, #2a2a2a); 
                                            color: #e9edef; padding: 15px 20px; 
                                            border-radius: 5px 25px 25px 25px; 
                                            max-width: 85%; border-left: 4px solid; 
                                            border-image: linear-gradient(180deg, #8b5cf6, #06b6d4) 1; 
                                            word-wrap: break-word;">
                                    <div style="white-space: pre-wrap;">{clean_res}</div>
                                </div>
                            </div>
                            """, unsafe_allow_html=True)
                            
                            last_render_time = current_time
                
                if res_text:
                    clean_res = clean_text(res_text)
                    response_container.markdown(f"""
                    <div style="display: flex; justify-content: flex-start; margin-bottom: 20px;">
                        {ai_avatar_html}
                        <div style="background: linear-gradient(135deg, #1a1a1a, #2a2a2a); 
                                    color: #e9edef; padding: 15px 20px; 
                                    border-radius: 5px 25px 25px 25px; 
                                    max-width: 85%; border-left: 4px solid; 
                                    border-image: linear-gradient(180deg, #8b5cf6, #06b6d4) 1; 
                                    word-wrap: break-word;">
                            <div style="white-space: pre-wrap;">{clean_res}</div>
                        </div>
                    </div>
                    """, unsafe_allow_html=True)
                    
                    res = res_text
        
        # ========== GROQ MODELS ==========
        elif engine_type == "Groq":
            messages = [{"role": "system", "content": system_prompt}]
            for m in st.session_state.messages[:-1]:
                if m.get("type") != "image":
                    messages.append({"role": m["role"], "content": m["content"]})
            messages.append({"role": "user", "content": final_user_msg})
            
            response_container = st.empty()
            res_text = ""
            
            last_render_time = time.time()
            RENDER_INTERVAL = 0.1
            
            ai_avatar_html = "<div style='width: 38px; height: 38px; border-radius: 50%; background: linear-gradient(135deg, #8b5cf6, #06b6d4); display: flex; align-items: center; justify-content: center; margin-right: 12px; border: 2px solid #06b6d4; font-size: 20px;'>🤖</div>"
            
            stream = client_groq.chat.completions.create(
                model=engine_config["model"],
                messages=messages,
                temperature=0.7,
                max_tokens=1024,
                stream=True
            )
            
            for chunk in stream:
                if chunk.choices[0].delta.content:
                    res_text += chunk.choices[0].delta.content
                    current_time = time.time()
                    
                    if current_time - last_render_time >= RENDER_INTERVAL:
                        clean_res = clean_text(res_text)
                        
                        response_container.markdown(f"""
                        <div style="display: flex; justify-content: flex-start; margin-bottom: 20px;">
                            {ai_avatar_html}
                            <div style="background: linear-gradient(135deg, #1a1a1a, #2a2a2a); 
                                        color: #e9edef; padding: 15px 20px; 
                                        border-radius: 5px 25px 25px 25px; 
                                        max-width: 85%; border-left: 4px solid; 
                                        border-image: linear-gradient(180deg, #8b5cf6, #06b6d4) 1; 
                                        word-wrap: break-word;">
                                <div style="white-space: pre-wrap;">{clean_res}</div>
                            </div>
                        </div>
                        """, unsafe_allow_html=True)
                        
                        last_render_time = current_time
            
            clean_res = clean_text(res_text)
            response_container.markdown(f"""
            <div style="display: flex; justify-content: flex-start; margin-bottom: 20px;">
                {ai_avatar_html}
                <div style="background: linear-gradient(135deg, #1a1a1a, #2a2a2a); 
                            color: #e9edef; padding: 15px 20px; 
                            border-radius: 5px 25px 25px 25px; 
                            max-width: 85%; border-left: 4px solid; 
                            border-image: linear-gradient(180deg, #8b5cf6, #06b6d4) 1; 
                            word-wrap: break-word;">
                    <div style="white-space: pre-wrap;">{clean_res}</div>
                </div>
            </div>
            """, unsafe_allow_html=True)
            
            res = res_text
        
        # ========== POLLINATIONS AI ==========
        elif engine_type == "Pollinations":
            # Image generation doesn't use file content, only the prompt
            encoded_prompt = urllib.parse.quote(user_msg)
            image_url = f"{POLLINATIONS_API}{encoded_prompt}"
            
            img_response = requests.get(image_url)
            img = Image.open(io.BytesIO(img_response.content))
            
            st.session_state.messages.append({"role": "assistant", "type": "image", "content": img})
            
            # Increment rate limit
            increment_rate_limit(st.session_state.current_user, selected_model, engine_config["limit"])
            
            if st.session_state.current_session_key:
                st.session_state.all_chats[st.session_state.current_session_key] = st.session_state.messages.copy()
            save_history_to_db(st.session_state.current_user, st.session_state.all_chats)
            st.rerun()
        
        if res:
            st.session_state.messages.append({"role": "assistant", "content": res})
            
            # Increment rate limit
            increment_rate_limit(st.session_state.current_user, selected_model, engine_config["limit"])
            
            if st.session_state.current_session_key:
                st.session_state.all_chats[st.session_state.current_session_key] = st.session_state.messages.copy()
            save_history_to_db(st.session_state.current_user, st.session_state.all_chats)
            st.rerun()
    
    except Exception as e:
        st.error(f"❌ Error bro: {str(e)}")
        error_msg = f"Sorry bro, ada error: {str(e)} 😰"
        st.session_state.messages.append({"role": "assistant", "content": error_msg})
        if st.session_state.current_session_key:
            st.session_state.all_chats[st.session_state.current_session_key] = st.session_state.messages.copy()
        save_history_to_db(st.session_state.current_user, st.session_state.all_chats)
        st.rerun()
